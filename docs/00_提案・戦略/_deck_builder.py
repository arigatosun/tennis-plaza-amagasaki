# -*- coding: utf-8 -*-
"""テニスプラザ尼崎 通し提案書(表紙+GA4現状報告→3プラン)PPTX生成。design.md準拠。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
GREEN=RGBColor(0x0F,0x3D,0x2E); GREEN2=RGBColor(0x26,0x48,0x31); GOLD=RGBColor(0xE8,0xA1,0x00)
CRIT=RGBColor(0xC0,0x39,0x2B); GOOD=RGBColor(0x1E,0x84,0x49); WARN=RGBColor(0xC7,0x7C,0x18)
INK=RGBColor(0x1F,0x2A,0x26); MUT=RGBColor(0x5D,0x6B,0x64); WARM=RGBColor(0xF7,0xF6,0xF2)
LINE=RGBColor(0xE3,0xE1,0xD8); WHITE=RGBColor(0xFF,0xFF,0xFF); CARD=RGBColor(0xFF,0xFF,0xFF)
LYEL=RGBColor(0xFB,0xF3,0xE6); LGREEN=RGBColor(0xED,0xF5,0xF0); LRED=RGBColor(0xFB,0xEE,0xEC)
GREY=RGBColor(0xBF,0xBF,0xBF)
FONT="Meiryo"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=13.333,7.5
BLANK=prs.slide_layouts[6]
_page=[0]

def _ea(run):
    rPr=run._r.get_or_add_rPr()
    for tag in ('a:ea','a:cs'):
        e=rPr.find(qn(tag))
        if e is None:
            e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',FONT)

def _noshadow(shp):
    try: shp.shadow.inherit=False
    except Exception: pass

def slide(bg=WARM):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); _noshadow(r)
    r.shadow.inherit=False
    return s

def tb(s,l,t,w,h,lines,anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {t,size,bold,color,align,space_after}"""
    box=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2);tf.margin_right=Pt(2);tf.margin_top=Pt(1);tf.margin_bottom=Pt(1)
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=ln.get('align',PP_ALIGN.LEFT)
        if ln.get('space_after') is not None: p.space_after=Pt(ln['space_after'])
        if ln.get('space_before') is not None: p.space_before=Pt(ln['space_before'])
        parts=ln['t'] if isinstance(ln['t'],list) else [(ln['t'],ln.get('color',INK),ln.get('bold',False))]
        for (txt,col,bold) in parts:
            r=p.add_run(); r.text=txt; r.font.size=Pt(ln.get('size',14)); r.font.bold=bold
            r.font.color.rgb=col; r.font.name=FONT; _ea(r)
    return box

def box(s,l,t,w,h,fill=CARD,line=LINE,line_w=1.0,radius=True,shadow=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False
    return shp

def shp_text(shp,lines,anchor=MSO_ANCHOR.MIDDLE):
    tf=shp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(8);tf.margin_right=Pt(8);tf.margin_top=Pt(5);tf.margin_bottom=Pt(5)
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=ln.get('align',PP_ALIGN.CENTER)
        if ln.get('space_after') is not None: p.space_after=Pt(ln['space_after'])
        parts=ln['t'] if isinstance(ln['t'],list) else [(ln['t'],ln.get('color',INK),ln.get('bold',False))]
        for (txt,col,bold) in parts:
            r=p.add_run(); r.text=txt; r.font.size=Pt(ln.get('size',13)); r.font.bold=bold
            r.font.color.rgb=col; r.font.name=FONT; _ea(r)

def title(s,txt,kicker=None):
    if kicker:
        tb(s,0.62,0.34,11,0.32,[{'t':kicker,'size':11,'bold':True,'color':GOLD}])
    tb(s,0.62,0.6,12.1,0.7,[{'t':txt,'size':27,'bold':True,'color':GREEN}])
    u=box(s,0.66,1.33,2.2,0.05,fill=GOLD,line=None,radius=False)
    footer(s)

def footer(s):
    _page[0]+=1
    tb(s,0.62,7.12,7,0.3,[{'t':'テニスプラザ尼崎 ご提案 ／ アリガトサン','size':9,'color':MUT}])
    tb(s,12.0,7.12,0.9,0.3,[{'t':str(_page[0]),'size':10,'color':MUT,'align':PP_ALIGN.RIGHT}])

def chevrons(s,l,t,items,active=None,w=2.7,h=0.85,gap=0.18,fill_on=GREEN,fill_off=RGBColor(0xDDE,0xDD,0xD7) if False else RGBColor(0xDD,0xDD,0xD7)):
    x=l
    for i,(lab,sub) in enumerate(items):
        c=s.shapes.add_shape(MSO_SHAPE.CHEVRON,Inches(x),Inches(t),Inches(w),Inches(h))
        on=(active is None or i<=active)
        c.fill.solid(); c.fill.fore_color.rgb=(GREEN if (active is not None and i==active) else (GREEN2 if on else RGBColor(0xD6,0xDb,0xD8)))
        c.line.fill.background(); c.shadow.inherit=False
        col=WHITE if on else MUT
        shp_text(c,[{'t':lab,'size':13,'bold':True,'color':col}]+([{'t':sub,'size':9,'color':(RGBColor(0xCF,0xE0,0xD8) if on else MUT)}] if sub else []))
        x+=w-0.35+gap

# ============ 1. 表紙 ============
s=slide(GREEN)
box(s,0,0,SW,SH,fill=GREEN,line=None,radius=False)
# accent bar
box(s,0,0,0.28,SH,fill=GOLD,line=None,radius=False)
tb(s,1.0,1.7,11,0.4,[{'t':'ARIGATOSAN ／ ご提案','size':13,'bold':True,'color':GOLD}])
tb(s,1.0,2.25,11.3,2.2,[
    {'t':'テニスプラザ尼崎さま','size':30,'bold':True,'color':WHITE,'space_after':6},
    {'t':'事業成長のご提案','size':40,'bold':True,'color':WHITE,'space_after':10},
    {'t':[('— 採用を起点に、',WHITE,False),('採れる→現場が回る→集客→西宮再建',GOLD,True),('まで —',WHITE,False)],'size':17,'color':WHITE},
])
box(s,1.0,5.5,11.0,0.02,fill=RGBColor(0x3a,0x5a,0x4c),line=None,radius=False)
tb(s,1.0,5.65,11,1.1,[
    {'t':'ご提案先：テニスプラザ尼崎 山下社長 ／ 近江さま','size':13,'color':RGBColor(0xE6,0xEF,0xEA),'space_after':3},
    {'t':'ご提案：アリガトサン（中田・中村） ／ 2026年6月','size':13,'color':RGBColor(0xCF,0xE0,0xD8),'space_after':3},
    {'t':'デザイン × マーケティング × 営業 × AIエンジニアリング × 経営','size':12,'color':GOLD},
])

# ============ 2. 本日のゴール ============
s=slide(); title(s,'本日のゴール と 進め方','AGENDA')
box(s,0.62,1.55,12.1,1.15,fill=LGREEN,line=GREEN,line_w=1.5)
tb(s,0.95,1.62,11.5,1.0,[
    {'t':[('本日のゴール：',GREEN,True),('「まず“穴を塞いで数えられる状態”を作る一歩（プランB）」のご判断をいただくこと。',INK,True)],'size':16,'space_after':2,'anchor':None},
    {'t':'※その先どこまで伸ばすか（プランC）は、数字が見えてから一緒に決められます。今日すべてを決める必要はありません。','size':12,'color':MUT},
],anchor=MSO_ANCHOR.MIDDLE)
ag=[('1. 現状報告（GA4の事実）','15分'),('2. 現状維持のコスト','5分'),('3. これからの方向性','10分'),
    ('4. 3つのプランご提案','15分'),('5. ご懸念への回答','5分'),('6. 次の一歩','10分')]
x=0.62;y=3.05;w=3.9;h=0.9
for i,(a,b) in enumerate(ag):
    col=i%3; row=i//3
    bx=box(s,0.62+col*4.0,3.05+row*1.05,3.8,0.92,fill=CARD,line=LINE)
    shp_text(bx,[{'t':a,'size':13,'bold':True,'color':GREEN},{'t':b,'size':11,'color':GOLD,'bold':True}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.62,5.5,12,0.4,[{'t':'（想定60分。今日いちばん大事なのは前半の「事実」より後半の「で、どうするか」です）','size':11,'color':MUT}])

# ============ 3. セクション扉: 現状報告 ============
def divider(num,title_txt,sub):
    s=slide(GREEN); box(s,0,0,SW,SH,fill=GREEN,line=None,radius=False)
    box(s,0,3.3,SW,0.04,fill=RGBColor(0x3a,0x5a,0x4c),line=None,radius=False)
    tb(s,1.1,2.2,11,0.7,[{'t':num,'size':15,'bold':True,'color':GOLD}])
    tb(s,1.1,2.75,11.2,1.2,[{'t':title_txt,'size':34,'bold':True,'color':WHITE}])
    tb(s,1.1,4.0,11.2,0.7,[{'t':sub,'size':15,'color':RGBColor(0xCF,0xE0,0xD8)}])
    footer(s); return s
divider('SECTION 1','現状報告 ─ まず、事実から','御社のGoogleアナリティクス（GA4）の実数を分析しました。')

# ============ 4. 結論+蛇口/バケツ/レジ/店員 ============
s=slide(); title(s,'結論：お客さんは来ている。でも「レジ」が無い。','現状をひとことで')
box(s,0.62,1.5,12.1,0.95,fill=GREEN,line=None)
shp_text(box(s,0.62,1.5,12.1,0.95,fill=GREEN,line=None),
         [{'t':'「集客が弱い」のではありません。土壌は良い。問題は “数えられていない” ことと “武器が繋がっていない” ことです。','size':16,'bold':True,'color':WHITE}],anchor=MSO_ANCHOR.MIDDLE)
# diagram row
dy=3.1; dh=2.0
labels=[("蛇口",GREEN,"SNS・検索\n（強い）","◎ 月1,650訪問\nリール18.5万再生"),
        ("バケツ",CRIT,"ホームページ\n（穴あき）","✕ 来た人を\n取りこぼし"),
        ("レジ",GREY,"計測\n（無い）","✕ 何件取れたか\n不明"),
        ("店員",WARN,"スタッフ\n（不足）","△ 近江さんが\n現場に張り付き")]
bw=2.7; gap=0.45; x0=0.95
for i,(name,col,sub,note) in enumerate(labels):
    x=x0+i*(bw+gap)
    b=box(s,x,dy,bw,dh,fill=WHITE,line=col,line_w=2.0)
    shp_text(b,[{'t':name,'size':20,'bold':True,'color':col,'space_after':2},
                {'t':sub.replace('\n','　'),'size':12,'bold':True,'color':INK,'space_after':3},
                {'t':note,'size':11,'color':MUT}],anchor=MSO_ANCHOR.MIDDLE)
    if i<3:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+bw+0.02),Inches(dy+dh/2-0.18),Inches(gap-0.04),Inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb=GOLD; a.line.fill.background(); a.shadow.inherit=False
tb(s,0.62,5.45,12.1,0.9,[
 {'t':[('まず ',INK,False),('バケツの穴を塞ぎ（HP改善）',CRIT,True),('、',INK,False),('レジを置き（計測）',GREEN,True),('、',INK,False),('店員を増やす（採用）',WARN,True),('。これが今日のご提案の全体像です。',INK,False)],'size':14}])

# ============ 5. GA4で見つけた問題 ①〜⑥ ============
s=slide(); title(s,'GA4で見つけた、6つの事実','現状報告｜すべて実データ＋スクショ証跡あり')
probs=[("①","成果(予約/問合せ/電話/入会)の計測がゼロ","レジが無い。何件取れているか不明",CRIT),
       ("②","Instagram→HPが月25訪問（全体の1.5%）","最強の武器がHPに繋がっていない",CRIT),
       ("③","有料SNS広告が滞在22秒で離脱","広告費の無駄打ち（出稿元も不明）",CRIT),
       ("④","Search Console未連携","どんな検索語で来ているか見えない",WARN),
       ("⑤","リンク切れ(404)が月275回","来た人を取りこぼしている",WARN),
       ("⑥","モバイル68%・月1,650訪問・検索の質76%","スマホ中心。土壌そのものは良い",GOOD)]
cw=5.95;ch=1.45;x0=0.62;y0=1.55
for i,(n,a,b,col) in enumerate(probs):
    col_i=i%2;row=i//2
    x=x0+col_i*6.15; y=y0+row*1.6
    c=box(s,x,y,cw,ch,fill=CARD,line=LINE)
    bdg=box(s,x+0.18,y+0.42,0.62,0.62,fill=col,line=None)
    shp_text(bdg,[{'t':n,'size':18,'bold':True,'color':WHITE}])
    tb(s,x+1.0,y+0.18,cw-1.15,ch-0.3,[
        {'t':a,'size':13,'bold':True,'color':INK,'space_after':2},
        {'t':b,'size':11.5,'color':col,'bold':True}],anchor=MSO_ANCHOR.MIDDLE)

# ============ 6. Instagram断絶 ============
s=slide(); title(s,'いちばんの“もったいない”：強い武器が繋がっていない','現状報告｜②の深掘り')
# left big
b=box(s,0.8,2.0,5.4,3.2,fill=LGREEN,line=GREEN,line_w=1.5)
shp_text(b,[{'t':'Instagram','size':16,'bold':True,'color':GREEN,'space_after':4},
            {'t':'18.5万','size':54,'bold':True,'color':GREEN,'space_after':0},
            {'t':'リール最大再生 ／ フォロワー893','size':12,'color':MUT}],anchor=MSO_ANCHOR.MIDDLE)
a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(6.35),Inches(3.35),Inches(0.7),Inches(0.5))
a.fill.solid();a.fill.fore_color.rgb=GOLD;a.line.fill.background();a.shadow.inherit=False
b2=box(s,7.15,2.0,5.4,3.2,fill=LRED,line=CRIT,line_w=1.5)
shp_text(b2,[{'t':'ホームページ到達','size':16,'bold':True,'color':CRIT,'space_after':4},
            {'t':'月25','size':54,'bold':True,'color':CRIT,'space_after':0},
            {'t':'訪問（全体のわずか1.5%）','size':12,'color':MUT}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.4,11.8,0.9,[{'t':[('これだけ見られているのに、HPにも予約にも繋がっていない。',CRIT,True),('　プロフィールの行き先を整えるだけで、ここは大きく伸びます。',INK,False)],'size':14}])

# ============ 7. CV計測ゼロ ============
s=slide(); title(s,'いま、成果を「数えていない」＝経営判断ができない','現状報告｜①の深掘り')
tb(s,0.62,1.55,12,0.6,[{'t':'GA4に記録されているのは自動の7項目だけ。予約・問い合わせ・電話・入会の「成果」が一つも設定されていません。','size':14,'color':INK}])
left=box(s,0.8,2.4,5.5,2.7,fill=CARD,line=LINE)
shp_text(left,[{'t':'記録されている','size':12,'color':MUT,'space_after':4},
   {'t':'ページ閲覧・スクロール 等','size':14,'bold':True,'color':INK,'space_after':6},
   {'t':'＝「人が来た」までは分かる','size':12,'color':MUT}],anchor=MSO_ANCHOR.MIDDLE)
right=box(s,6.6,2.4,5.9,2.7,fill=LRED,line=CRIT,line_w=1.5)
shp_text(right,[{'t':'記録されていない','size':12,'color':CRIT,'space_after':4},
   {'t':'予約・問い合わせ・電話・入会','size':15,'bold':True,'color':CRIT,'space_after':6},
   {'t':'＝「何件 成果が出たか」が0件のまま','size':12,'color':INK,'bold':True}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.35,11.8,0.8,[{'t':[('例えるなら「レジの無いお店」。',CRIT,True),('どの施策が効いたか分からず、広告も改善も“勘”になります。まずここを直します。',INK,False)],'size':14}])

# ============ 8. 土壌は良い ============
s=slide(); title(s,'でも、土壌は良い。だから“伸びしろ”です','現状報告｜⑥のポジティブ面')
stats=[("月1,650","訪問／月（決して少なくない）",GREEN),
       ("76%","検索で来た人の“しっかり見た”率",GREEN),
       ("68%","スマホからのアクセス",GREEN)]
for i,(big,sub,col) in enumerate(stats):
    x=0.8+i*4.0
    b=box(s,x,2.2,3.7,2.6,fill=LGREEN,line=GREEN,line_w=1.2)
    shp_text(b,[{'t':big,'size':44,'bold':True,'color':GREEN,'space_after':2},{'t':sub,'size':12,'color':INK}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.2,11.8,0.9,[{'t':[('集客力が無いのではありません。',INK,False),('来てくれた人を取りこぼさず、数える仕組みを作るだけ',GREEN,True),('で成果は変わります。',INK,False)],'size':14}])

# ============ 9. 現状維持のコスト ============
s=slide(); title(s,'「今のままでいい」の、見えていない値段','現状報告｜現状維持のコスト試算')
box(s,0.62,1.5,12.1,0.85,fill=LRED,line=CRIT,line_w=1.2)
shp_text(box(s,0.62,1.5,12.1,0.85,fill=LRED,line=CRIT),
   [{'t':'計測ゼロ・導線の穴のままだと、「取れていたはずの成果」が毎月見えないまま消えています。','size':15,'bold':True,'color':CRIT}],anchor=MSO_ANCHOR.MIDDLE)
# illustrative calc
rows=[("入会1件の価値（月謝1万円×平均在籍2年）","約24万円【仮置き】"),
      ("導線改善で体験申込が月+3件 → 入会 月+1.2件（入会率40%・仮）","月+1.2件"),
      ("年間の機会損失（24万円 × 1.2件 × 12ヶ月）","約345万円／年【仮置き】")]
y=2.6
for i,(a,b) in enumerate(rows):
    bx=box(s,0.8,y,11.7,0.78,fill=(LYEL if i<2 else GREEN),line=LINE if i<2 else None)
    shp_text(bx,[{'t':[(a+'　→　',(INK if i<2 else WHITE),i==2),(b,(CRIT if i<2 else WHITE),True)],'size':14}],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.92
tb(s,0.8,5.5,11.8,0.95,[
  {'t':'※数字はすべて仮置きの試算です。実際の会員数・月謝・体験→入会率はMTGで伺い、その場で「逆算シート」に入れて確定値を出します。','size':11,'color':MUT,'space_after':2},
  {'t':[('要点：直すコストより、',INK,False),('放置するコストの方が大きい',CRIT,True),('可能性が高い、ということです。',INK,False)],'size':13}])

# ============ 10. セクション扉: 方向性 ============
divider('SECTION 2','これからの方向性','穴を塞ぎ、数え、人を増やす。その正しい順番。')

# ============ 11. なぜ計測の次に採用か(接続) ============
s=slide(); title(s,'なぜ、最初の一歩が「採用」なのか','方向性｜計測 と 採用 をつなぐ')
tb(s,0.62,1.5,12,0.55,[{'t':'「HPの穴を塞ぐ（計測）」と並んで、もう一つ最初に効くドミノが「採用」です。理由はシンプルです。','size':14}])
steps=[("近江さんが\n現場に張り付き","集客や改善を\n回す人がいない","だから\nまず採用","現場が回る→\n手離れ→集客→西宮")]
items=[("①近江さんが現場に張り付き","今は近江さん依存"),("②集客・改善を回す人がいない","施策が回らない"),
       ("③だから まず採用","最初のドミノ"),("④現場が回る→手離れ→集客→西宮","好循環へ")]
chevrons(s,0.7,2.6,items,active=2,w=3.05,h=1.0)
box(s,0.8,4.3,11.7,1.0,fill=LGREEN,line=GREEN,line_w=1.2)
shp_text(box(s,0.8,4.3,11.7,1.0,fill=LGREEN,line=GREEN),
  [{'t':[('レジを直しても、店員がいなければ店は回らない。',GREEN,True),('だから「計測」と「採用」を同時に、が正解です。',INK,False)],'size':15}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.5,11.8,0.7,[{'t':'※これは私たちの分析からのご提言です。御社のご事情に合わせて一緒に調整できます。','size':11,'color':MUT}])

# ============ 12. 全体ロードマップ ============
s=slide(); title(s,'尼崎で型を作り、西宮へ。全体ロードマップ','方向性｜本日のスコープと将来')
phases=[("PHASE 1","採用","社員1+バイト2を採用"),("PHASE 1","近江 手離れ","現場の比重を下げる"),
        ("PHASE 2","集客強化","取りこぼしを止め伸ばす"),("PHASE 3","西宮 再建","尼崎の型を横展開")]
x0=0.7
for i,(p,a,b) in enumerate(phases):
    x=x0+i*3.05
    on = i<2
    c=s.shapes.add_shape(MSO_SHAPE.PENTAGON,Inches(x),Inches(2.5),Inches(3.0),Inches(1.5))
    c.fill.solid();c.fill.fore_color.rgb=(GREEN if on else RGBColor(0xD6,0xDB,0xD8));c.line.fill.background();c.shadow.inherit=False
    shp_text(c,[{'t':p,'size':10,'bold':True,'color':(GOLD if on else MUT)},
                {'t':a,'size':17,'bold':True,'color':(WHITE if on else MUT)},
                {'t':b,'size':9.5,'color':(RGBColor(0xCF,0xE0,0xD8) if on else MUT)}],anchor=MSO_ANCHOR.MIDDLE)
box(s,0.7,4.35,5.95,0.6,fill=GOLD,line=None)
shp_text(box(s,0.7,4.35,5.95,0.6,fill=GOLD,line=None),[{'t':'◀ 本日ご提案するスコープ','size':12,'bold':True,'color':GREEN}])
tb(s,0.8,5.4,11.8,0.9,[{'t':[('西宮再建は“3年後のゴール”として見据えますが、',INK,False),('本日決めるのは尼崎の第一歩（採用＋計測）だけ',GREEN,True),('です。まずここから。',INK,False)],'size':14}])

# ============ 13. 採用ペルソナ（近江さん定義） ============
s=slide(); title(s,'採りたいのは、こんなコーチ（近江さん定義）','方向性｜誰を採るか')
shp_text(box(s,0.62,1.45,12.1,0.8,fill=LYEL,line=GOLD,line_w=1.2),[{'t':'「レッスンのコマをたくさん持って、長く現場に立ち続けてくれるコーチ」','size':16,'bold':True,'color':GREEN}],anchor=MSO_ANCHOR.MIDDLE)
shp_text(box(s,0.62,2.45,5.9,2.0,fill=LGREEN,line=GREEN,line_w=1.3),[{'t':'◎ 採りたい','size':14,'bold':True,'color':GREEN,'space_after':4},
  {'t':'・コマを安定して回す現場コーチ','size':12.5,'color':INK,'space_after':2,'align':PP_ALIGN.LEFT},
  {'t':'・同じことを長く続けられる','size':12.5,'color':INK,'space_after':2,'align':PP_ALIGN.LEFT},
  {'t':'・若くて元気・一定の給料と休みで満足','size':12.5,'color':INK,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
shp_text(box(s,6.82,2.45,5.9,2.0,fill=LRED,line=CRIT,line_w=1.3),[{'t':'✕ あえて求めない','size':14,'bold':True,'color':CRIT,'space_after':4},
  {'t':'・「営業もできます」系の器用貧乏','size':12.5,'color':INK,'space_after':2,'align':PP_ALIGN.LEFT},
  {'t':'・向上心が強すぎ、すぐ次を求める人','size':12.5,'color':INK,'space_after':2,'align':PP_ALIGN.LEFT},
  {'t':'・ベテラン／昭和世代','size':12.5,'color':INK,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
types=[("① テニス一筋の若手","社員候補"),("② 社会人から戻った若手","社員候補"),("③ 副業・かけもち","バイト/委託")]
for i,(t,sub) in enumerate(types):
    bx=box(s,0.62+i*4.05,4.65,3.85,0.78,fill=CARD,line=LINE)
    shp_text(bx,[{'t':[(t+'　',INK,True),(sub,GOLD,True)],'size':12.5}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.62,5.6,12.1,0.95,[
  {'t':[('条件：',GREEN,True),('社員1人＝年収280〜300万・週2.5コマ〜 ／ バイト2人＝月4.5万(交通費込)・各4コマ ／ 平成生まれ・元気重視。',INK,False)],'size':12.5,'space_after':2},
  {'t':'※近江さんの手書きメモ(2026-06-03)を整理。＝「私たちはちゃんと聞きました」を示す1枚。','size':11,'color':MUT}])

# ============ 13b. KPI設計：採用逆算(中村式) ============
s=slide(); title(s,'KPI設計：採用3名までを“数字”で逆算','方向性｜中村方式：何人集めれば何人採れるか')
tb(s,0.62,1.45,12.1,0.6,[{'t':[('採用 ',INK,False),('←',GOLD,True),(' 面接通過率 ',INK,False),('←',GOLD,True),(' 応募(LP閲覧×CVR) ',INK,False),('←',GOLD,True),(' 各チャネル流入。',INK,False),('この逆算で「どこから何人」まで落とします。',GREEN,True)],'size':13}])
chain=[("各チャネル流入 → LP閲覧","500",GREEN2,"CVR 10%"),("応募","50",GREEN,"面談化 44%"),("カジュアル面談","22",GREEN,"通過 50%"),("最終面接","9",WARN,"採用 33%"),("採用（社員1＋バイト2）","3",GOLD,"")]
y=2.25; fullw=8.6
for i,(lab,num,col,rate) in enumerate(chain):
    w=fullw*(1-i*0.13)
    bx=box(s,0.8,y,w,0.66,fill=col,line=None)
    shp_text(bx,[{'t':[(lab+'　',WHITE,True),(num+('名' if i>0 else ''),WHITE,True)],'size':13}],anchor=MSO_ANCHOR.MIDDLE)
    if rate: tb(s,0.8+w+0.18,y+0.08,3.4,0.5,[{'t':'▼ '+rate,'size':11,'bold':True,'color':MUT}])
    y+=0.82
tb(s,0.8,6.05,11.9,0.7,[{'t':'※CVR10%・各通過率は中村KPI＋業界相場の仮置き。実数を入れれば「逆算シート」が自動計算します。','size':11,'color':MUT}])

# ============ 13c. どこから何人 ============
s=slide(); title(s,'どこから、何人 連れてくるか','方向性｜チャネル別の具体数字（すべて無料媒体）')
dd=[["チャネル（無料）","流入(LP閲覧)","応募","この数字の作り方"],
 ["Instagram／リール＋公式LINE","175 人","18 人","リール約1.8万再生→プロフ→LP(遷移約1%)。最大の伸びしろ"],
 ["無料求人エンジン(engage→Indeed/しごと検索)","150 人","15 人","1回登録で複数媒体へ自動掲載＋無料スカウト"],
 ["リファラル＋地域(ハローワーク/ジモティー)","125 人","12 人","生徒・OB・体育会＋地域無料。テニス好きに直撃"],
 ["チラシ／オフライン(高校・大学→QR)","50 人","5 人","Z世代・新卒枠。近江さんが配布"],
 ["合計","500 人","50 人","→ 最終面接 9名 → 採用 3名"]]
gt=s.shapes.add_table(len(dd),4,Inches(0.62),Inches(1.7),Inches(12.1),Inches(3.6)).table
gt.columns[0].width=Inches(4.3);gt.columns[1].width=Inches(1.9);gt.columns[2].width=Inches(1.3);gt.columns[3].width=Inches(4.6)
for ri,row in enumerate(dd):
    gt.rows[ri].height=Inches(0.55)
    for ci,val in enumerate(row):
        cell=gt.cell(ri,ci);cell.margin_left=Pt(6);cell.margin_right=Pt(4);cell.margin_top=Pt(2);cell.margin_bottom=Pt(2);cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.text_frame.word_wrap=True;p=cell.text_frame.paragraphs[0]
        p.alignment=PP_ALIGN.CENTER if 1<=ci<=2 else PP_ALIGN.LEFT
        r=p.add_run();r.text=val;r.font.name=FONT;_ea(r)
        header=(ri==0);last=(ri==len(dd)-1)
        r.font.size=Pt(12 if not header else 12.5);r.font.bold=header or last or ci==2
        if header: r.font.color.rgb=WHITE;cell.fill.solid();cell.fill.fore_color.rgb=GREEN
        elif last: r.font.color.rgb=(GREEN if ci<3 else INK);cell.fill.solid();cell.fill.fore_color.rgb=LGREEN
        else:
            r.font.color.rgb=(CRIT if ci==2 else INK);cell.fill.solid();cell.fill.fore_color.rgb=(WARM if ri%2 else WHITE)
tb(s,0.62,5.55,12.1,0.95,[{'t':[('ペルソナ＝「テニス好き・長く続けたい」若手（社員）＋副業層（バイト）。',INK,True)],'size':11.5,'space_after':2},{'t':'数字は仮置き。実数（インスタ再生・HP流入・LINE友だち）を入れれば「逆算シート」が自動更新します。','size':11,'color':MUT}])

# ============ 14. 媒体: 人材紹介を使わない ============
s=slide(); title(s,'採用に“消える費用”は使いません','方向性｜採用媒体の方針')
l=box(s,0.8,1.8,5.7,3.0,fill=LRED,line=CRIT,line_w=1.5)
shp_text(l,[{'t':'人材紹介（doda/リクルート）','size':14,'bold':True,'color':CRIT,'space_after':4},
   {'t':'約90〜100万円','size':40,'bold':True,'color':CRIT,'space_after':2},
   {'t':'社員1人あたり（年収280-300万の30-35%）','size':11,'color':MUT,'space_after':4},
   {'t':'✕ 採るたびに消える','size':12,'bold':True,'color':CRIT}],anchor=MSO_ANCHOR.MIDDLE)
r=box(s,6.85,1.8,5.7,3.0,fill=LGREEN,line=GREEN,line_w=1.5)
shp_text(r,[{'t':'私たちの設計（無料媒体）','size':14,'bold':True,'color':GREEN,'space_after':4},
   {'t':'ほぼ ¥0','size':40,'bold':True,'color':GREEN,'space_after':2},
   {'t':'engage一括配信／Googleしごと検索／リファラル／ハローワーク／ジモティー','size':10.5,'color':MUT,'space_after':4},
   {'t':'◎ 仕組みが御社に残る（唯一の有料補強＝スポキャリ¥55,000）','size':11,'bold':True,'color':GREEN}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.15,11.8,0.9,[{'t':[('この“使わずに済む採用コスト”が、',INK,False),('そのまま制作・改善・仕組みづくりの原資',GOLD,True),('になります。',INK,False)],'size':15}])

# ============ 15. セクション扉: ご提案 ============
divider('SECTION 3','ご提案 ─ 3つのプラン','「どこまでやるか」を、御社が選べる形にしました。')

# ============ 16. 3プラン全体像 ============
s=slide(); title(s,'3つのプラン','ご提案｜全体像')
plans=[("PLAN A","採用の入口","まず、人を採れる状態に","40〜65万円","単発・目安",GREEN2,None),
       ("PLAN B","バケツの穴を塞ぐ","穴を塞ぎ、レジを置く","90〜120万円","単発・目安",GOLD,"推奨"),
       ("PLAN C","数えて磨いて伴走","直すだけで終わらせない","初期120〜160万＋月15万","初期＋伴走・目安",GREEN,"本命")]
cw=3.9;x0=0.62
for i,(p,nm,catch,price,term,col,badge) in enumerate(plans):
    x=x0+i*4.05
    c=box(s,x,1.7,3.85,4.4,fill=CARD,line=col,line_w=(2.4 if badge else 1.0))
    if badge:
        bd=box(s,x+1.1,1.5,1.65,0.42,fill=col,line=None)
        shp_text(bd,[{'t':badge,'size':12,'bold':True,'color':(GREEN if col==GOLD else WHITE)}])
    tb(s,x+0.25,1.95,3.4,3.9,[
        {'t':p,'size':12,'bold':True,'color':GOLD,'space_after':1,'align':PP_ALIGN.CENTER},
        {'t':nm,'size':19,'bold':True,'color':GREEN,'space_after':1,'align':PP_ALIGN.CENTER},
        {'t':catch,'size':11,'color':MUT,'space_after':8,'align':PP_ALIGN.CENTER},
        {'t':price,'size':25,'bold':True,'color':INK,'space_after':0,'align':PP_ALIGN.CENTER},
        {'t':term+'（税別）','size':10,'color':MUT,'space_after':8,'align':PP_ALIGN.CENTER},
    ])
abc=[("採用LP・しごと検索・採用導線","＋HPの穴を塞ぐ・計測・LINE・AIスカウト","＋HP刷新・ブランド型化・運用伴走")]
notes=["採用の一点突破","推奨：土台を一度の工事で","本命：手離れ→西宮まで伴走"]
for i in range(3):
    x=x0+i*4.05
    nb=box(s,x,4.65,3.85,1.35,fill=(LYEL if i==1 else LGREEN if i==2 else WARM),line=LINE)
    shp_text(nb,[{'t':abc[0][i],'size':11,'bold':True,'color':INK,'space_after':3},{'t':notes[i],'size':10.5,'color':MUT}],anchor=MSO_ANCHOR.MIDDLE)

tb(s,0.62,6.18,12.1,0.7,[{'t':'※金額は目安です。ご予算・承認フローに合わせて調整します。初回をオーナー承認不要の範囲で小さく始め、実績を見て伸ばす進め方も可能です。','size':11,'color':MUT}])

# ============ 17. プランA ============
def plan_detail(code,name,catch,price,term,col,badge,bullets,foot,foot_col=MUT):
    s=slide(); title(s,f'{code}：{name}',f'ご提案｜{catch}')
    pb=box(s,9.0,1.5,3.6,1.25,fill=col,line=None)
    shp_text(pb,[{'t':code+('　'+badge if badge else ''),'size':13,'bold':True,'color':(GREEN if col==GOLD else WHITE),'space_after':2},
                 {'t':price,'size':22,'bold':True,'color':(GREEN if col==GOLD else WHITE)},
                 {'t':term+'（税別）','size':10,'color':(GREEN if col==GOLD else RGBColor(0xE6,0xEF,0xEA))}],anchor=MSO_ANCHOR.MIDDLE)
    y=1.6
    for (txt,add) in bullets:
        bx=box(s,0.7,y,8.0,0.62,fill=(LYEL if add else CARD),line=LINE)
        shp_text(bx,[{'t':[(('＋ ' if add else '✓ '),(GOLD if add else GOOD),True),(txt,INK,False)],'size':12.5,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
        y+=0.7
    fb=box(s,0.7,y+0.05,11.9,0.95,fill=LGREEN,line=GREEN,line_w=1.2)
    shp_text(fb,[{'t':foot,'size':13,'bold':True,'color':GREEN}],anchor=MSO_ANCHOR.MIDDLE)
    return s
plan_detail("PLAN A","採用の入口","まず、人を採れる状態をつくる","40〜65万円","単発・目安",GREEN2,None,
    [("採用ペルソナ設計（近江さん定義＝若く長く続く現場コーチ）",False),
     ("採用LPの磨き込み・本番稼働化（既存資産を活用）",False),
     ("Googleしごと検索の有効化（無料掲載）",False),
     ("Instagram採用導線・採用バナー・OGP",False),
     ("採用CVの計測ポイント設計",False)],
    "想定応募：30日で 10〜20件（しごと検索＋無料媒体の露出見込み・仮置き）／ まず小さく承認を通したい方へ")

# ============ 18. プランB ============
plan_detail("PLAN B","バケツの穴を塞ぐ","蛇口は強い。穴を塞ぎ、レジを置く","90〜120万円","単発・目安",GOLD,"推奨",
    [("プランA の内容すべて",False),
     ("CV計測の実装＝「レジを置く」（予約/問合せ/電話/入会など主要6種）",True),
     ("Instagram→HP断絶の解消・有料広告の正体特定（UTM）",True),
     ("Search Console連携 ／ 404リンク切れ修正 ／ スマホ最適化",True),
     ("無料求人媒体の面展開 ＋ 公式LINE開設",True),
     ("AIで応募者への声かけを下書き自動化（最終送信は人）",True),
     ("KPIダッシュボード（社長が一目で見られる1枚）",True)],
    "★初期設定後はダッシュボードを見るだけ。運用の手間は最小（近江さん・スタッフでも回せる手順込み）")

# ============ 19. プランB 投資回収 ============
s=slide(); title(s,'プランB は、何ヶ月で回収できるか','ご提案｜投資回収シミュレーション')
box(s,0.62,1.5,12.1,0.7,fill=LGREEN,line=GREEN)
shp_text(box(s,0.62,1.5,12.1,0.7,fill=LGREEN,line=GREEN),
  [{'t':'投資 約90〜120万円。入会1件の価値を約24万円【仮置き】とすると──','size':14,'bold':True,'color':GREEN}],anchor=MSO_ANCHOR.MIDDLE)
cases=[("保守ケース","追加入会 月+1件","24万円/月","約5ヶ月で回収",GREEN2),
       ("標準ケース","追加入会 月+2件","48万円/月","約3ヶ月で回収",GREEN)]
for i,(t,a,b,c,col) in enumerate(cases):
    x=0.8+i*6.05
    bx=box(s,x,2.5,5.85,1.9,fill=CARD,line=col,line_w=1.5)
    shp_text(bx,[{'t':t,'size':13,'bold':True,'color':col,'space_after':2},
                 {'t':a+'　→　'+b,'size':13,'color':INK,'space_after':3},
                 {'t':c,'size':20,'bold':True,'color':col}],anchor=MSO_ANCHOR.MIDDLE)
box(s,0.8,4.65,11.85,0.95,fill=LYEL,line=GOLD,line_w=1.2)
shp_text(box(s,0.8,4.65,11.85,0.95,fill=LYEL,line=GOLD),
  [{'t':[('さらに、',INK,False),('社員1人を人材紹介で採ると約90〜100万円',CRIT,True),('。それを使わないだけで、Bの費用はほぼ相殺されます。',INK,False)],'size':14}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.75,11.8,0.5,[{'t':'※数字は仮置き。実数（会員数・月謝・入会率）をMTGで伺い、逆算シートで確定値をその場でお出しします。','size':11,'color':MUT}])

# ============ 20. プランC ============
plan_detail("PLAN C","数えて、磨いて、伴走する","直すだけで終わらせない。回し続ける","初期120〜160万＋月15万","初期＋伴走・目安",GREEN,"本命",
    [("プランB の内容すべて",False),
     ("HP全体のビジュアル刷新・情報設計",True),
     ("ブランドフレームワーク（Instagram投稿/リール型/運用マニュアル）",True),
     ("AIロープレ（保護者対応・入会クロージングの練習）",True),
     ("LINE自動応答 ／ SEOコラム（AI下書き＋監修）",True),
     ("月次レポート＋戦略MTG／採用KPI進捗管理（小改修込み）",True)],
    "月15万の中身＝毎月の“お約束”：月次レポート＋戦略MTG／施策の実行／改善／3ヶ月で中間レビュー（縮小・停止OK）")

# ============ 21. C 残る/消える 対比 ============
s=slide(); title(s,'同じご予算をかけるなら「残るか・消えるか」','ご提案｜プランC の価値')
l=box(s,0.8,1.8,5.7,3.4,fill=LRED,line=CRIT,line_w=1.5)
shp_text(l,[{'t':'人材紹介にお金をかけても','size':15,'bold':True,'color':CRIT,'space_after':6},
   {'t':'・人は採れる','size':13,'color':INK,'space_after':3},
   {'t':'・でも、それで終わり','size':13,'color':INK,'space_after':3},
   {'t':'・資産は残らない','size':13,'color':INK,'space_after':3},
   {'t':'・来年もまた同じ費用','size':13,'bold':True,'color':CRIT}],anchor=MSO_ANCHOR.MIDDLE)
r=box(s,6.85,1.8,5.7,3.4,fill=LGREEN,line=GREEN,line_w=1.5)
shp_text(r,[{'t':'同じご予算をプランCにかけると','size':14,'bold':True,'color':GREEN,'space_after':6},
   {'t':'・人も採れる（無料媒体）','size':13,'color':INK,'space_after':3},
   {'t':'・HP/ブランド/計測基盤が残る','size':13,'color':INK,'space_after':3},
   {'t':'・運用を回す力が社内に残る','size':13,'color':INK,'space_after':3},
   {'t':'・来年以降の採用・集客コストが下がる','size':13,'bold':True,'color':GREEN}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.8,5.4,11.8,0.7,[{'t':[('同じ金額なら、',INK,False),('“消える費用”より“残る資産”。',GREEN,True),('これがプランCの考え方です。',INK,False)],'size':15}])

# ============ 22. 比較表 ============
s=slide(); title(s,'3プラン 比較表','ご提案｜何が付くか一目で')
data=[["含まれるもの","A","B","C"],
 ["採用LP・しごと検索・採用導線","○","○","○"],
 ["① CV計測（レジ設置・主要6種）","−","○","◎"],
 ["② Instagram→HP導線・③広告UTM","−","○","◎"],
 ["④ Search Console ／ ⑤404 ／ ⑥スマホ最適化","−","○","◎"],
 ["無料求人媒体の面展開・公式LINE","−","○","◎"],
 ["AI半自動スカウト ／ KPIダッシュボード","−","○","◎"],
 ["HP全体刷新・ブランド型化・運用マニュアル","−","−","○"],
 ["AIロープレ／LINE自動応答／SEOコラム","−","−","○"],
 ["月次レポート＋戦略MTG（伴走・小改修込み）","−","−","○"],
 ["初期費用（目安・税別）","40〜65万","90〜120万","120〜160万"],
 ["月額（税別）","−","−","15万/月"]]
rows=len(data);cols=4
gt=s.shapes.add_table(rows,cols,Inches(0.7),Inches(1.55),Inches(11.9),Inches(5.1)).table
gt.columns[0].width=Inches(7.1);gt.columns[1].width=Inches(1.6);gt.columns[2].width=Inches(1.6);gt.columns[3].width=Inches(1.6)
for ri,row in enumerate(data):
    gt.rows[ri].height=Inches(0.4)
    for ci,val in enumerate(row):
        cell=gt.cell(ri,ci); cell.margin_left=Pt(6);cell.margin_right=Pt(4);cell.margin_top=Pt(1);cell.margin_bottom=Pt(1)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
        p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=val; r.font.name=FONT; _ea(r)
        header=(ri==0); last=(ri>=rows-2)
        r.font.size=Pt(12 if not header else 12.5); r.font.bold=header or last or ci==0
        if header:
            r.font.color.rgb=WHITE; cell.fill.solid(); cell.fill.fore_color.rgb=GREEN
        else:
            if ci==2: cell.fill.solid(); cell.fill.fore_color.rgb=LYEL
            elif last: cell.fill.solid(); cell.fill.fore_color.rgb=LGREEN
            else: cell.fill.solid(); cell.fill.fore_color.rgb=(WARM if ri%2 else WHITE)
            if val=="◎": r.font.color.rgb=GOOD
            elif val=="○": r.font.color.rgb=GREEN
            elif val=="−": r.font.color.rgb=GREY
            else: r.font.color.rgb=INK

# ============ 23. なぜアリガトサンか ============
s=slide(); title(s,'なぜ、私たち（アリガトサン）か','ご提案｜他社との違い')
cols3=[("制作会社","作って終わり","・きれいなHPは作る\n・でも“その後”は無い\n・成果は測らない",GREY,LRED),
       ("人材紹介","人だけ・毎回消える","・人は連れてくる\n・成功報酬30-35%\n・仕組みは残らない",CRIT,LRED),
       ("アリガトサン","数えて・回して・残す","・デザイン×マーケ×営業\n　×AI×経営の総合力\n・計測して改善し続ける\n・仕組みが御社に残る",GREEN,LGREEN)]
for i,(t,sub,body,col,bg) in enumerate(cols3):
    x=0.62+i*4.05
    c=box(s,x,1.7,3.85,3.7,fill=bg,line=col,line_w=(2.2 if i==2 else 1.0))
    shp_text(c,[{'t':t,'size':17,'bold':True,'color':col,'space_after':2},
                {'t':sub,'size':12,'bold':True,'color':INK,'space_after':6},
                {'t':body,'size':12,'color':INK}],anchor=MSO_ANCHOR.TOP)
tb(s,0.62,5.65,12,0.8,[{'t':[('御社は既に良い土壌と強いSNSをお持ちです。それを',INK,False),('“成果に変えて、御社に残す”',GREEN,True),('のが私たちの役割です。',INK,False)],'size':14}])

# ============ 24. 想定されるご懸念 Q&A ============
s=slide(); title(s,'想定されるご懸念に、先にお答えします','ご提案｜Q&A')
qa=[("「今のままでも回ってる」","“見えていないだけ”です。計測を入れると、取りこぼしが数字で見えます（前述の試算）。"),
    ("「高いのでは？」","人材紹介なら一度で消える同額が、Cなら資産として残ります。回収も数字で示せます。"),
    ("「効果が出なかったら？」","初月で“何件取れているか”を可視化。Cは3ヶ月で中間レビュー（縮小・停止OK）。"),
    ("「他社と比べたい」","制作会社＝作って終わり／紹介＝人だけ。私たちは数えて回して“残す”点が違います。")]
for i,(q,a) in enumerate(qa):
    y=1.65+i*1.25
    qb=box(s,0.7,y,3.7,1.05,fill=LRED,line=CRIT,line_w=1.0)
    shp_text(qb,[{'t':'Q. '+q,'size':13,'bold':True,'color':CRIT}],anchor=MSO_ANCHOR.MIDDLE)
    ab=box(s,4.6,y,7.95,1.05,fill=CARD,line=LINE)
    shp_text(ab,[{'t':[('A. ',GREEN,True),(a,INK,False)],'size':12.5,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)

# ============ 25. おすすめ B→C ============
s=slide(); title(s,'おすすめは「B」。ゴールは「C」。でも今日はBだけで充分。','ご提案｜進め方')
flow=[("① B でレジを置く","“何件取れているか”が初めて数字で見える"),
      ("② 数字が次の問いを生む","「型化したら何件？」「広告を運用に回したら？」"),
      ("③ 数字を見てCへ","必要だと分かってから伸ばす（C前提でBを買う必要なし）"),
      ("④ 西宮へ","回り出した実績が、上の方への最大の説得材料に")]
for i,(a,b) in enumerate(flow):
    y=1.7+i*1.05
    n=box(s,0.7,y,4.7,0.92,fill=(LYEL if i==0 else LGREEN),line=(GOLD if i==0 else GREEN),line_w=1.2)
    shp_text(n,[{'t':a,'size':14,'bold':True,'color':(GOLD if i==0 else GREEN)}],anchor=MSO_ANCHOR.MIDDLE)
    ab=box(s,5.6,y,6.95,0.92,fill=CARD,line=LINE)
    shp_text(ab,[{'t':b,'size':12.5,'color':INK,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.7,6.05,11.9,0.7,[{'t':[('押し売りはしません。',INK,False),('数字が次のステップを連れてきます。',GREEN,True)],'size':14}])

# ============ 26. ご契約後の流れ＋お約束 ============
s=slide(); title(s,'ご契約後の流れと、初月のお約束','ご提案｜進め方')
tl=[("1週","キックオフ・権限確認・ペルソナ確定"),("2-3週","CV計測（レジ）実装・採用LP・LINE開設"),
    ("3-4週","404修正・導線・無料媒体展開／中間確認"),("5-8週","スマホ最適化・ダッシュボード納品")]
x0=0.7
for i,(w,d) in enumerate(tl):
    x=x0+i*3.05
    hd=box(s,x,1.7,2.9,0.55,fill=GREEN,line=None)
    shp_text(hd,[{'t':w+'目','size':13,'bold':True,'color':WHITE}])
    bd=box(s,x,2.3,2.9,1.5,fill=CARD,line=LINE)
    shp_text(bd,[{'t':d,'size':11.5,'color':INK}],anchor=MSO_ANCHOR.MIDDLE)
    if i<3:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+2.9),Inches(2.85),Inches(0.16),Inches(0.4))
        a.fill.solid();a.fill.fore_color.rgb=GOLD;a.line.fill.background();a.shadow.inherit=False
box(s,0.7,4.2,11.9,1.0,fill=LGREEN,line=GREEN,line_w=1.4)
shp_text(box(s,0.7,4.2,11.9,1.0,fill=LGREEN,line=GREEN),
  [{'t':[('初月のお約束：',GREEN,True),('「何件取れているか」が数字で見える状態にします（見えなければ計測部分は作り直し）。',INK,False)],'size':14}],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.7,5.4,11.9,0.7,[{'t':'プランCは3ヶ月時点で中間レビュー。指標が動かなければ縮小・停止できます（リスクは私たちが取ります）。','size':12,'color':MUT}])

# ============ 27. クロージング ============
s=slide(GREEN); box(s,0,0,SW,SH,fill=GREEN,line=None,radius=False); box(s,0,0,0.28,SH,fill=GOLD,line=None,radius=False)
tb(s,1.0,1.4,11,0.4,[{'t':'NEXT STEP','size':13,'bold':True,'color':GOLD}])
tb(s,1.0,2.0,11.3,1.6,[
   {'t':'まず穴を塞いで、“何件取れているか”が','size':26,'bold':True,'color':WHITE,'space_after':2},
   {'t':'見える状態をつくりましょう。','size':26,'bold':True,'color':WHITE}])
tb(s,1.0,3.9,11.3,1.3,[
   {'t':[('本日 ',RGBColor(0xE6,0xEF,0xEA),False),('プランB',GOLD,True),(' でご判断いただければ、来週から計測設置に着手し、',RGBColor(0xE6,0xEF,0xEA),False)],'size':15,'space_after':3},
   {'t':[('初月末に「最初の数字」',GOLD,True),('をお見せします。その先どこまで伸ばすかは、数字を見て一緒に。',RGBColor(0xE6,0xEF,0xEA),False)],'size':15}])
box(s,1.0,5.6,5.5,0.9,fill=GOLD,line=None)
shp_text(box(s,1.0,5.6,5.5,0.9,fill=GOLD,line=None),[{'t':'まずは、Bで第一歩を。','size':18,'bold':True,'color':GREEN}])
footer(s)

# ============ 28. 上申用1枚サマリー ============
s=slide(); title(s,'【ご参考】上の方へお持ちいただく1枚','APPENDIX｜決裁用サマリー')
tb(s,0.62,1.4,12,0.4,[{'t':'山下社長が、上のオーナー様にそのまま見せられる要点1枚です。','size':12,'color':MUT}])
sm=[("現状の問題","成果(予約/問合せ/入会)の計測ゼロ＋強いSNSがHPに繋がっていない（月25訪問/1.5%）。404が月275回。"),
    ("放置のコスト","機会損失 年 約345万円【仮置き】。広告費も効果不明のまま流出。"),
    ("打ち手","採用（無料媒体中心・人材紹介＝社員1人約90-100万を回避）＋ HPの穴を塞ぎ計測を入れる。"),
    ("推奨プラン","B＝90〜120万円（単発・目安）。回収 約3〜5ヶ月【仮置き】。本命Cは初期120〜160万＋月15万。"),
    ("残る資産","採用の仕組み・HP改善・計測基盤・運用力が社内に残る（来年以降のコストが下がる）。"),
    ("出口条件","初月で成果を可視化。Cは3ヶ月で中間レビュー（縮小・停止可）。")]
y=2.0
for i,(k,v) in enumerate(sm):
    kb=box(s,0.7,y,2.7,0.72,fill=GREEN,line=None)
    shp_text(kb,[{'t':k,'size':12.5,'bold':True,'color':WHITE}])
    vb=box(s,3.55,y,9.0,0.72,fill=(WARM if i%2 else WHITE),line=LINE)
    shp_text(vb,[{'t':v,'size':12,'color':INK,'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.8

# ============ 29. 裏表紙 ============
s=slide(GREEN); box(s,0,0,SW,SH,fill=GREEN,line=None,radius=False)
tb(s,1.0,3.0,11.3,1.6,[
   {'t':'ありがとうございました。','size':30,'bold':True,'color':WHITE,'space_after':6},
   {'t':'テニスプラザ尼崎さまの“次の一歩”を、全力でご一緒します。','size':15,'color':RGBColor(0xCF,0xE0,0xD8)}])
tb(s,1.0,5.4,11,0.5,[{'t':'アリガトサン（中田・中村） ／ デザイン × マーケ × 営業 × AI × 経営','size':13,'color':GOLD}])
footer(s)

import os
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"tennis_plaza_proposal_deck.pptx")
prs.save(out)
print("OK slides=",len(prs.slides._sldIdLst),"->",out)
